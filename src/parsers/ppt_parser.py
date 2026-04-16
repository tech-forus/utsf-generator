"""
PowerPoint Parser — PPTX / PPT
================================
Extracts text, tables, and logistics data from PowerPoint presentations.

Often used by transporters for:
- Rate card presentations (slide tables with zone rates)
- Company profile decks (company info, contact details)
- Network coverage maps (text with zone/city lists)

Pipeline:
  1. python-pptx to extract text + tables from all slides
  2. PDF-style charge extraction on combined text
  3. Zone matrix detection on slide tables
  4. Company info extraction (GST, phone, email)
"""

import os
import re
from typing import Dict, List, Any, Optional
from parsers.base_parser import BaseParser


# Reuse the same patterns as the PDF parser
_CHARGE_PATTERNS = [
    ("fuel",           r"fuel\s*(?:surcharge|%|percent)[^\d]*(\d+(?:\.\d+)?)"),
    ("docketCharges",  r"(?:docket|doc(?:ument)?)\s*(?:charges?|fee)[^\d]*(?:rs\.?\s*)?(\d+(?:\.\d+)?)"),
    ("minCharges",     r"min(?:imum)?\s*(?:charges?|freight)[^\d]*(?:rs\.?\s*)?(\d+(?:\.\d+)?)"),
    ("minWeight",      r"min(?:imum)?\s*(?:chargeable\s*)?weight[^\d]*(\d+(?:\.\d+)?)\s*kg"),
    ("greenTax",       r"green\s*(?:tax|cess|levy)[^\d]*(\d+(?:\.\d+)?)"),
    ("rovCharges_v",   r"r\.?o\.?v\.?[^\d]*(\d+(?:\.\d+)?)\s*%"),
    ("odaCharges_f",   r"o\.?d\.?a\.?[^\d]*(?:rs\.?\s*)?(\d+(?:\.\d+)?)\s*(?:per\s*(?:shipment|consignment))?"),
    ("divisor",        r"(?:volumetric\s*divisor|kfactor|k\s*factor)[^\d]*(\d+)"),
]

_COMPANY_PATTERNS = {
    "gstNo":        r'\b(\d{2}[A-Z]{5}\d{4}[A-Z]\d[ZYX]\d)\b',
    "panNo":        r'\bPAN\s*[:\-]?\s*([A-Z]{5}\d{4}[A-Z])\b',
    "contactPhone": r'\b((?:\+?91[-\s]?)?\d{10})\b',
    "contactEmail": r'\b([\w.+-]+@[\w-]+\.[\w.]+)\b',
    "website":      r'\b((?:https?://)?(?:www\.)?[\w-]+\.(?:com|in|co\.in|net|org)(?:/[\w/.-]*)?)\b',
}


class PPTParser(BaseParser):
    SUPPORTED_EXTENSIONS = [".pptx", ".ppt"]

    def parse(self, file_path: str) -> Dict[str, Any]:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pptx":
            return self._parse_pptx(file_path)
        else:
            # .ppt (old binary format) — try converting via pptx or return limited parse
            return self._parse_ppt_fallback(file_path)

    # ── PPTX (python-pptx) ────────────────────────────────────────────────────

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
        except ImportError:
            print("[PPTParser] python-pptx not installed — pip install python-pptx")
            return {"text": "", "tables": [], "data": {}}

        prs = Presentation(file_path)
        all_text_parts = []
        all_tables: List[List[List[str]]] = []
        sheets: Dict[str, List[List[str]]] = {}

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = []

            for shape in slide.shapes:
                # Text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            slide_text_parts.append(line)

                # Tables
                if shape.has_table:
                    tbl = shape.table
                    rows: List[List[str]] = []
                    for row in tbl.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        # Deduplicate adjacent merged cells
                        deduped = []
                        for c in cells:
                            if not deduped or c != deduped[-1]:
                                deduped.append(c)
                        if any(deduped):
                            rows.append(deduped)

                    if rows:
                        name = f"Slide{slide_idx + 1}_Table{len(all_tables) + 1}"
                        all_tables.append(rows)
                        sheets[name] = rows

            if slide_text_parts:
                all_text_parts.append(f"[Slide {slide_idx + 1}]")
                all_text_parts.extend(slide_text_parts)

        full_text = "\n".join(all_text_parts)

        return self._extract_data(full_text, all_tables, sheets)

    def _parse_ppt_fallback(self, file_path: str) -> Dict[str, Any]:
        """For .ppt files: try LibreOffice conversion, else limited extraction."""
        print(f"[PPTParser] .ppt binary format — limited support. "
              "Convert to .pptx in PowerPoint for best results.")
        return {"text": "", "tables": [], "data": {}, "sheets": {}}

    # ── Data extraction ────────────────────────────────────────────────────────

    def _extract_data(
        self,
        text: str,
        tables: List[List[List[str]]],
        sheets: Dict[str, List[List[str]]] = None,
    ) -> Dict[str, Any]:
        text_lower = text.lower()
        data: Dict[str, Any] = {}

        # Company info
        company = self._extract_company_info(text)
        if company:
            data["company_details"] = company

        # Charges from text
        charges = self._extract_charges(text)
        if charges:
            data["charges"] = charges

        # Zone matrix from tables — delegate to base parser via sheets dict
        if sheets:
            data["sheets"] = sheets
        if tables:
            data["tables"] = tables

        data["text"] = text
        return data

    def _extract_company_info(self, text: str) -> Dict[str, Any]:
        info = {}
        for field, pattern in _COMPANY_PATTERNS.items():
            m = re.search(pattern, text, re.IGNORECASE)
            if m:
                info[field] = m.group(1).strip()
        # Try to extract company name from first non-empty line
        for line in text.splitlines():
            line = line.strip()
            if len(line) > 5 and not any(c.isdigit() for c in line[:5]):
                # Remove slide marker
                if not line.startswith("[Slide"):
                    info.setdefault("companyName", line)
                    break
        return info

    def _extract_charges(self, text: str) -> Dict[str, Any]:
        charges = {}
        text_lower = text.lower()
        for field, pattern in _CHARGE_PATTERNS:
            m = re.search(pattern, text_lower)
            if m:
                try:
                    val = float(m.group(1))
                    if field.endswith("_v"):
                        base = field[:-2]
                        charges.setdefault(base, {})["v"] = val
                    elif field.endswith("_f"):
                        base = field[:-2]
                        charges.setdefault(base, {})["f"] = val
                    else:
                        charges[field] = val
                except ValueError:
                    pass
        return charges
